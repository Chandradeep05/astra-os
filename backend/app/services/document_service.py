import os
import uuid
import json
import aiohttp
import difflib
from typing import List, Optional, Dict, Any
from pydantic import BaseModel
import pypdf
from docx import Document
from pptx import Presentation
import logging
import re
from app.core.config import settings
from app.services.vector_service import vector_service

# Import multi-modal tools
try:
    from app.tools.vision_tool import run_vision_ocr
    from app.tools.audio_tool import run_audio_transcription
except ImportError:
    run_vision_ocr = None
    run_audio_transcription = None

try:
    from sentence_transformers import CrossEncoder
except ImportError:
    CrossEncoder = None

logger = logging.getLogger(__name__)

# ASTRA-FIX: Module-level retrieval confidence thresholds
# ChromaDB default = L2 distance. nomic-embed-text L2 distances:
#   Very relevant:  0.3 - 0.8
#   Somewhat:       0.8 - 1.2
#   Irrelevant:     1.2+
# We convert to a 0-1 similarity scale using: sim = 1 / (1 + distance)
# This maps: dist=0 → sim=1.0, dist=0.5 → sim=0.67, dist=1.0 → sim=0.50, dist=2.0 → sim=0.33
# nomic-embed-text produces L2 distances of 400-600 for .docx content.
# Formula: sim = 1/(1+L2). At L2=477 → sim=0.002. Keyword boost adds up to +0.10.
# Thresholds must match the ACTUAL distribution, not theoretical cosine ranges.
HIGH_CONF = 0.05   # Chunk with decent keyword overlap (~3+ matching words)
LOW_CONF  = 0.01   # Chunk with at least some keyword overlap (~1 word)


# ── Filename extraction for source-level filtering ────────────────────────────
# CRITICAL: Do NOT allow spaces in filename pattern — it greedily matches
# the entire query prefix as a filename (e.g. "Explain everything from FILE.docx"
# was being detected as filename "Explain everything from FILE.docx").
# Underscores and hyphens are fine: circuit_system_solutions.docx, ai-notes.pdf
_FILENAME_PATTERN = re.compile(
    r'([\w\-]+\.(?:docx|pdf|txt|xlsx|pptx|csv|md|json))',
    re.IGNORECASE,
)

def _normalize_filename(name: str) -> str:
    """Normalize filename for fuzzy matching: strip separators, lowercase.
    'AI_Notes.pdf' → 'ainotes.pdf', 'ai-notes.PDF' → 'ainotes.pdf'"""
    base, ext = os.path.splitext(name)
    normalized = re.sub(r'[\s\-_]+', '', base).lower()
    return normalized + ext.lower()

def _extract_filename_from_query(query: str):
    """Extract a filename like 'circuit_system_solutions.docx' from a user query."""
    match = _FILENAME_PATTERN.search(query)
    return match.group(1) if match else None

# Pronoun/reference patterns — "summarize it", "the first document", "this file"
_PRONOUN_REF_PATTERN = re.compile(
    r'\b(summarize it|explain it|explain that|what does it say|what does it mention|'
    r'read it|open it|the same document|the same file|that document|that file|'
    r'this document|this file|the document|the file|'
    r'the (first|second|third|last|next|previous) (document|file|pdf|upload))\b',
    re.IGNORECASE,
)

# Fix #10: Max characters per chunk to prevent token explosion
_MAX_CHUNK_CHARS = 2000


class DocumentService:
    def __init__(self):
        self.db_url = settings.DATABASE_URL
        self.db_type = "postgres" if "postgresql" in self.db_url.lower() else "sqlite"
        self.ollama_url = settings.OLLAMA_BASE_URL
        self.model = "nomic-embed-text:latest" # Standard high-quality local embedding model
        self._http_session: Optional[aiohttp.ClientSession] = None
        self._last_queried_source: Optional[str] = None  # Fix #2: track last successfully queried source for pronoun resolution
        self._last_uploaded_source: Optional[str] = None
        self._last_uploaded_at: Optional[float] = None  # timestamp

    async def _get_session(self) -> aiohttp.ClientSession:
        """Get or create a shared HTTP session for embedding calls."""
        if self._http_session is None or self._http_session.closed:
            self._http_session = aiohttp.ClientSession()
        return self._http_session

    async def close(self):
        """Clean up the HTTP session on shutdown."""
        if self._http_session and not self._http_session.closed:
            await self._http_session.close()

    async def get_embedding(self, text: str) -> List[float]:
        """Generate embedding using local Ollama instance (Async).

        Retries up to 3 times with backoff to handle Ollama being busy
        serving the chat model (common on single-GPU systems).
        Supports both legacy /api/embeddings and newer /api/embed endpoints.
        """
        import asyncio as _asyncio
        # /api/embeddings (legacy), /api/embed (Ollama >=0.5), /api/embedding (typo guard)
        endpoints = ["/api/embeddings", "/api/embed", "/api/embedding"]
        max_retries = 3

        session = await self._get_session()

        for attempt in range(max_retries):
            for endpoint in endpoints:
                # /api/embed uses {"input": ...}, legacy uses {"prompt": ...}
                body_key = "input" if endpoint == "/api/embed" else "prompt"
                try:
                    async with session.post(
                        f"{self.ollama_url}{endpoint}",
                        json={"model": self.model, body_key: text},
                        timeout=aiohttp.ClientTimeout(total=60)
                    ) as response:
                        if response.status == 200:
                            data = await response.json()
                            # /api/embed returns {"embeddings": [...]}, legacy returns {"embedding": [...]}
                            embedding = data.get("embedding") or (data.get("embeddings", [None]) or [None])[0]
                            if not embedding:
                                logger.warning(f"Empty embedding from {endpoint}")
                                continue
                            return embedding
                        elif response.status == 404:
                            continue  # Try next endpoint
                        else:
                            error_text = await response.text()
                            logger.warning(f"Ollama {endpoint} returned {response.status}: {error_text}")
                except Exception as e:
                    logger.warning(f"Embedding attempt {attempt+1} on {endpoint}: {e}")
                    continue

            # Wait before retry — Ollama may be busy with the chat model
            if attempt < max_retries - 1:
                wait = 5 * (attempt + 1)
                logger.info(f"[EMBEDDING-RETRY] Waiting {wait}s before retry {attempt+2}/{max_retries}")
                await _asyncio.sleep(wait)

        logger.error(f"Embedding failed for model {self.model} after {max_retries} retries.")
        raise RuntimeError(f"Failed to generate embeddings after {max_retries} retries. Ensure Ollama is running and model '{self.model}' is installed.") 

    def _chunk_text(self, text: str, max_size: int = 1000, overlap: int = 150) -> List[str]:
        """Semantically chunks text by paragraphs and sentences with overlap."""
        import re
        chunks = []
        paragraphs = re.split(r'\n\n+', text.strip())
        current_chunk = ""

        def add_to_chunk(segment: str):
            nonlocal current_chunk
            if not current_chunk:
                current_chunk = segment
            elif len(current_chunk) + len(segment) + 1 <= max_size:
                current_chunk += " " + segment
            else:
                chunks.append(current_chunk.strip())
                # Create overlap without cutting words in half
                overlap_text = current_chunk[-overlap:]
                last_space = overlap_text.rfind(' ')
                if last_space != -1:
                    overlap_text = overlap_text[last_space+1:]
                new_chunk = (overlap_text.strip() + " " + segment).strip()
                if len(new_chunk) > max_size:
                    chunks.append(overlap_text.strip())
                    current_chunk = segment
                else:
                    current_chunk = new_chunk

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            if len(para) <= max_size:
                add_to_chunk(para)
            else:
                # Split paragraph into sentences safely
                sentences = re.split(r'(?<=[.!?])\s+', para)
                for sent in sentences:
                    if len(sent) <= max_size:
                        add_to_chunk(sent)
                    else:
                        # Fallback for massive sentences: split by words
                        words = sent.split(' ')
                        temp = ""
                        for word in words:
                            if len(temp) + len(word) + 1 <= max_size:
                                temp += (" " + word) if temp else word
                            else:
                                add_to_chunk(temp)
                                temp = word
                        if temp:
                            add_to_chunk(temp)
                            
        if current_chunk:
            chunks.append(current_chunk.strip())
            
        return [c for c in chunks if len(c) > 50]

    async def process_and_index_file(
        self,
        file_path: str,
        project_id: str,
        original_filename: str = None,  # ASTRA-FIX: real name before UUID rename
        file_id: str = None,            # ASTRA-FIX-2: unique ID for delete safety
    ) -> bool:
        """Extract text, chunk it, generate embeddings concurrently, and store in ChromaDB."""
        text = await self.process_file(file_path)
        if not text:
            return False

        # Semantic Chunking
        chunks = self._chunk_text(text)

        try:
            import asyncio
            collection_name = f"project_{project_id}"
            
            # Generate embeddings concurrently with a strict limit to avoid overloading Ollama
            semaphore = asyncio.Semaphore(1)
            async def limited_embedding(chunk_str):
                async with semaphore:
                    return await self.get_embedding(chunk_str)
            embeddings = await asyncio.gather(*[limited_embedding(chunk) for chunk in chunks])
            
            ids = [uuid.uuid4().hex for _ in chunks]
            file_name = (original_filename or os.path.basename(file_path)).lower()  # normalize to lowercase for consistent case-insensitive matching
            # ASTRA-FIX-2: Use explicit file_id (UUID) when provided.
            # Fallback to filename stem ONLY for backward compat with old callers.
            effective_file_id = file_id or os.path.splitext(file_name)[0]
            # Fix #3: Store normalized filename for fuzzy matching
            source_norm = _normalize_filename(file_name)
            metadatas = [
                {
                    "source": file_name,
                    "source_normalized": source_norm,
                    "file_id": effective_file_id,
                    "project_id": project_id,
                }
                for _ in chunks
            ]
                
            success = vector_service.add_documents(
                collection_name=collection_name,
                documents=chunks,
                ids=ids,
                metadatas=metadatas,
                embeddings=list(embeddings)
            )
            import time
            self._last_uploaded_source = file_name
            self._last_uploaded_at = time.time()

            return success
        except Exception as e:
            logger.error(f"Error indexing file {file_path}: {e}")
            return False

    async def search_similar(  # ASTRA-FIX
        self,
        query: str,
        limit: int = 5,               # ASTRA-FIX: project_id no longer positional
        file_id: str = None,          # ASTRA-FIX: optional specific-file filter
        project_id: str = None,       # ASTRA-FIX: now optional for backward compat
        min_similarity: float = None, # ASTRA-FIX: backward-compat override
        query_class: str = None,      # FIX-3: classifier output — skip ChromaDB unless RAG_QUERY
    ) -> dict:                        # ASTRA-FIX: always dict, never list/tuple
        """Confidence-tiered vector search in ChromaDB (Async).

        Returns:
            {
                "results": [{"content", "metadata", "similarity", "confidence"}],
                "confidence_level": "high" | "low" | "none"
            }
        """
        # FIX-3: Short-circuit — only RAG_QUERY ever touches ChromaDB.
        # All other classes (DIRECT_LLM, TOOL_CALL, etc.) return empty immediately.
        if query_class is not None and query_class != "RAG_QUERY":
            logger.debug(f"[RETRIEVAL-SKIP] query_class={query_class} — skipping ChromaDB")
            return {"results": [], "confidence_level": "none"}


        # ASTRA-FIX: Step 1 — Clean query for better retrieval
        # Strip trigger phrases that add noise to the embedding
        cleaned_query = self._clean_query_for_retrieval(query)  # ASTRA-FIX
        logger.info(f"[ASTRA-RETRIEVAL] Original: '{query[:80]}' → Cleaned: '{cleaned_query[:80]}'")  # ASTRA-FIX

        query_embedding = await self.get_embedding(cleaned_query)  # ASTRA-FIX: embed cleaned query

        try:
            # Auto-detect filename in query for source-level filtering
            detected_source = _extract_filename_from_query(query)

            # ASTRA-UPGRADE: Bug #3 — Ordinal resolution ("the second document")
            import re as _re_ord
            ordinal_match = _re_ord.search(
                r'\bthe (first|second|third|last|1st|2nd|3rd)\b', query, _re_ord.IGNORECASE
            )
            if ordinal_match and not detected_source and not file_id:
                ordinal_word = ordinal_match.group(1).lower()
                ordinal_map = {
                    "first": 0, "1st": 0,
                    "second": 1, "2nd": 1,
                    "third": 2, "3rd": 2,
                    "last": -1,
                }
                idx = ordinal_map.get(ordinal_word)
                if idx is not None:
                    docs_in_order = await self._get_documents_ordered_by_upload(project_id or "default")
                    if docs_in_order:
                        target = docs_in_order[idx] if idx != -1 else docs_in_order[-1]
                        detected_source = target["original_name"]
                        logger.info(f"[ASTRA-ORDINAL] Resolved '{ordinal_word}' → {detected_source}")

            # Fix #2: Pronoun resolution — "summarize it" → use last queried source
            # ASTRA-UPGRADE: Prefer recently uploaded file (within 5 minutes) over last queried
            if not detected_source and not file_id and _PRONOUN_REF_PATTERN.search(query):
                import time
                if (
                    self._last_uploaded_source
                    and self._last_uploaded_at
                    and (time.time() - self._last_uploaded_at) < 300
                ):
                    detected_source = self._last_uploaded_source
                    logger.info(f"[ASTRA-PRONOUN] Resolved to recently uploaded: {detected_source}")
                elif self._last_queried_source:
                    detected_source = self._last_queried_source
                    logger.info(f"[ASTRA-PRONOUN] Resolved to last queried: {detected_source}")

            # ── FUZZY FILENAME MATCHING (DO THIS WEEK) ───────────────────────────
            # Users often mistype filenames. If detected_source doesn't match any
            # exact file in the project, try fuzzy matching with a 0.8 threshold.
            if detected_source and not file_id:
                # Get all docs for this project
                available_docs = await self._get_documents_ordered_by_upload(project_id or "default")
                doc_names = [d["original_name"] for d in available_docs]
                
                # Check for exact match first (case-insensitive)
                exact_match = next((name for name in doc_names if name.lower() == detected_source.lower()), None)
                
                if not exact_match and doc_names:
                    # FIX: compare lowercase on both sides
                    detected_lower = detected_source.lower()
                    doc_names_lower = [n.lower() for n in doc_names]
                    matches_lower = difflib.get_close_matches(detected_lower, doc_names_lower, n=1, cutoff=0.75)
                    if matches_lower:
                        matched_idx = doc_names_lower.index(matches_lower[0])
                        logger.info(f"[ASTRA-FUZZY] Corrected '{detected_source}' -> '{doc_names[matched_idx]}'")
                        detected_source = doc_names[matched_idx]
                elif exact_match:
                    detected_source = exact_match
            # ───────────────────────────────────────────────────────────────────────

            # ───────────────────────────────────────────────────────────────────────



            # Build where clause from file_id and/or detected source filename
            # Fix #3: Try source_normalized for fuzzy matching, fall back to source for old chunks
            where = None
            _used_normalized = False  # Track which filter was used for fallback logic
            if file_id and detected_source:
                where = {"$and": [{"file_id": {"$eq": file_id}}, {"source_normalized": {"$eq": _normalize_filename(detected_source)}}]}
                _used_normalized = True
            elif file_id:
                where = {"file_id": {"$eq": file_id}}
            elif detected_source:
                detected_norm = _normalize_filename(detected_source)
                where = {"source_normalized": {"$eq": detected_norm}}
                _used_normalized = True
                logger.info(f"[ASTRA-FILTER] Auto-detected source filter (normalized): {detected_norm}")

            if where:
                logger.debug(f"[ASTRA-FILTER-APPLIED] where={where}")

            # Use collection name scoped to project when available
            collection_name = f"project_{project_id}" if project_id else "default"

            # Over-fetch for better recall, then rerank
            fetch_limit = max(limit * 2, 10)

            query_kwargs: Dict[str, Any] = {
                "collection_name": collection_name,
                "query_embeddings": query_embedding,
                "n_results": fetch_limit,
            }
            if where is not None:
                query_kwargs["where"] = where

            results = vector_service.query(**query_kwargs)

            # Backward compat: if source_normalized returned nothing, retry with source (lowercase)
            # Old chunks indexed before the normalization fix only have the "source" field
            if _used_normalized and (
                not results
                or not results.get("documents")
                or len(results["documents"]) == 0
                or len(results["documents"][0]) == 0
            ):
                logger.info("[ASTRA-COMPAT] source_normalized returned empty — retrying with source (lowercase)")
                fallback_source = detected_source.lower()
                if file_id:
                    where = {"$and": [{"file_id": {"$eq": file_id}}, {"source": {"$eq": fallback_source}}]}
                else:
                    where = {"source": {"$eq": fallback_source}}
                query_kwargs["where"] = where
                results = vector_service.query(**query_kwargs)

            # ASTRA-FIX: Robust empty-response guard
            if (
                not results
                or not results.get("documents")
                or len(results["documents"]) == 0
                or len(results["documents"][0]) == 0
            ):
                logger.debug("[ASTRA-EMPTY] ChromaDB returned no documents.")  
                return {
                    "results": [], 
                    "confidence_level": "none",
                    "guard_triggered": True
                }

            docs  = results["documents"][0]
            metas = results["metadatas"][0] if results.get("metadatas") else [{}] * len(docs)
            dists = (
                results["distances"][0]
                if results.get("distances") and results["distances"]
                else [0.0] * len(docs)
            )

            # ASTRA-FIX: Correct L2-safe similarity mapping
            # ChromaDB default distance = L2 (Euclidean), NOT cosine.
            # L2 range is [0, ∞), NOT [0, 1]. Using 1-dist gives NEGATIVE similarity for dist > 1.
            # Fix: use inverse mapping: similarity = 1 / (1 + distance)
            # This gives a clean 0.0–1.0 range regardless of distance metric.
            effective_low = min_similarity if min_similarity is not None else LOW_CONF

            # ASTRA-FIX: Step 4 — Initialize CrossEncoder lazily (Step 2)
            if not hasattr(self, "_cross_encoder"):
                if CrossEncoder is not None:
                    logger.info("[ASTRA-RETRIEVAL] Loading CrossEncoder: ms-marco-MiniLM-L-6-v2")
                    self._cross_encoder = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2', max_length=512)
                else:
                    self._cross_encoder = None

            query_keywords = set(cleaned_query.lower().split())

            # ── SCORING ──────────────────────────────────────────────────────────────────

            # Batch predict CrossEncoder scores if available
            ce_scores = None
            if self._cross_encoder is not None:
                try:
                    pairs = [[cleaned_query, doc] for doc in docs]
                    import asyncio
                    raw_scores = await asyncio.to_thread(self._cross_encoder.predict, pairs)
                    ce_scores = raw_scores
                    logger.info(
                        f"[ASTRA-RERANK] CrossEncoder: "
                        f"min={raw_scores.min():.3f}, max={raw_scores.max():.3f}, "
                        f"mean={raw_scores.mean():.3f}"
                    )
                except Exception as e:
                    logger.error(f"[ASTRA-RETRIEVAL] CrossEncoder prediction failed: {e}")
                    ce_scores = None

            all_scored: List[Dict[str, Any]] = []

            for i in range(len(docs)):
                meta     = metas[i] or {}
                distance = dists[i] if dists[i] is not None else 0.0
                vector_sim = round(1.0 / (1.0 + distance), 4)

                if ce_scores is not None:
                    # CrossEncoder raw score — higher = more relevant
                    # Typical range: -10 to +10
                    ce_score = float(ce_scores[i])
                    final_sim = ce_score
                    keyword_boost = 0.0
                else:
                    # Fallback: vector similarity + keyword overlap boost
                    chunk_words = set(docs[i].lower().split())
                    query_words = set(cleaned_query.lower().split())
                    overlap = len(query_words & chunk_words)
                    keyword_boost = min(overlap * 0.02, 0.10)
                    final_sim = round(min(vector_sim + keyword_boost, 1.0), 4)

                all_scored.append({
                    "content":       docs[i],
                    "metadata":      meta,
                    "similarity":    final_sim,
                    "vector_sim":    vector_sim,
                    "keyword_boost": keyword_boost,
                    "l2_distance":   round(distance, 4),
                })

            # Sort ALL results by final similarity DESC
            all_scored.sort(key=lambda x: x["similarity"], reverse=True)

            # ── ADAPTIVE THRESHOLDS ──────────────────────────────────────────────────────
            # Thresholds depend on scoring mode (cross-encoder vs vector+keyword)

            if ce_scores is not None:
                # Cross-encoder mode: scores typically in [-10, +10]
                score_range = max(r["similarity"] for r in all_scored) - min(r["similarity"] for r in all_scored)
                effective_high = 2.0   # High relevance: score > 2
                effective_low  = 0.0   # Low relevance: score > 0 (positive = relevant at all)
                logger.info(f"[ASTRA-THRESH] CE mode: high={effective_high}, low={effective_low}")
            else:
                # Vector+keyword mode: use fixed thresholds
                effective_high = HIGH_CONF   # 0.05
                effective_low  = effective_low if 'effective_low' in dir() else LOW_CONF  # 0.01 or overridden

            # If user explicitly requested a document, lower threshold but don't disable
            if where is not None:
                effective_low = max(effective_low * 0.5, 0.0)
                logger.info(f"[ASTRA-RETRIEVAL] Explicit doc filter active. Lowered threshold to {effective_low:.4f}")

            # ── TIER CLASSIFICATION ──────────────────────────────────────────────────────

            high_results: List[Dict[str, Any]] = []
            low_results:  List[Dict[str, Any]] = []

            for r in all_scored:
                # Fix #10: Truncate oversized chunks to prevent token explosion
                content = r["content"]
                if len(content) > _MAX_CHUNK_CHARS:
                    content = content[:_MAX_CHUNK_CHARS] + "..."
                result_item = {
                    "content":    content,
                    "metadata":   r["metadata"],
                    "similarity": r["similarity"],
                }

                if r["similarity"] >= effective_high:
                    result_item["confidence"] = "high"
                    high_results.append(result_item)
                elif r["similarity"] >= effective_low:
                    result_item["confidence"] = "low"
                    low_results.append(result_item)
                else:
                    logger.debug(
                        f"[ASTRA-FILTER] Discarded chunk similarity={r['similarity']:.3f} "
                        f"source={r['metadata'].get('source', 'unknown')}"
                    )

            # ── TIER SELECTION ───────────────────────────────────────────────────────────

            if high_results:
                # Track last successful source for pronoun resolution
                # Use detected_source if available, otherwise extract from top chunk metadata
                if detected_source:
                    self._last_queried_source = detected_source.lower()
                else:
                    top_source = high_results[0].get("metadata", {}).get("source")
                    if top_source:
                        self._last_queried_source = top_source
                        logger.info(f"[ASTRA-TRACK] Tracked source from result: {top_source}")
                logger.info(f"[ASTRA-RESULT] Returning {min(len(high_results), limit)} HIGH confidence chunks")
                return {
                    "results": high_results[:limit],
                    "confidence_level": "high",
                    "guard_triggered": False
                }
            elif low_results:
                # Track last successful source for pronoun resolution
                if detected_source:
                    self._last_queried_source = detected_source.lower()
                else:
                    top_source = low_results[0].get("metadata", {}).get("source")
                    if top_source:
                        self._last_queried_source = top_source
                
                selected = low_results[:min(limit, 3)]
                logger.info(f"[ASTRA-RESULT] Returning {len(selected)} LOW confidence chunks (fallback)")
                return {
                    "results": selected,
                    "confidence_level": "low",
                    "guard_triggered": False
                }
            else:
                # ASTRA-UPGRADE: Bug #C — Strict Source Filtering
                # If a source was explicitly requested, DO NOT fall back to unfiltered search.
                # Returning results from OTHER documents causes "Context Leakage" (Doc Bleed).
                if where is not None and (detected_source or file_id):
                    logger.info(f"[ASTRA-STRICT] Source filter '{detected_source or file_id}' returned nothing. Refusing cross-doc fallback.")
                    return {
                        "results": [],
                        "confidence_level": "none",
                        "guard_triggered": True
                    }

                logger.info("[ASTRA-GUARD] No chunks passed similarity threshold. Triggering hallucination guard.")
                return {
                    "results": [],
                    "confidence_level": "none",
                    "guard_triggered": True
                }

        except Exception as e:
            logger.error(f"Error searching vector store: {e}")
            return {"results": [], "confidence_level": "none", "guard_triggered": True}

    def _clean_query_for_retrieval(self, query: str) -> str:
        """Extract semantic core from query for embedding.

        Strategy: remove conversational wrapper, keep topic words.
        'What does the document say about Z-transform?' → 'Z-transform'
        'Explain all concepts from CIRCUIT_SYSTEM_SOLUTIONS.DOCX' → 'concepts'
        'summarize this document' → '' → fallback to topic anchor
        """
        cleaned = query

        # Step 1: Strip filename — it's used for metadata filter, not embedding
        detected = _extract_filename_from_query(query)
        if detected:
            cleaned = cleaned.replace(detected, "")
            cleaned = cleaned.replace(detected.upper(), "")
            cleaned = cleaned.replace(detected.lower(), "")

        # Step 2: Strip doc-reference wrapper phrases
        NOISE_PHRASES = [
            r"according to (my|this|the|that) (document[s]?|file[s]?|note[s]?|upload[s]?)",
            r"(from|in|inside) (my|this|the|that) (document[s]?|file[s]?|note[s]?|upload[s]?)",
            r"(my|this|the|that) (document[s]?|file[s]?|note[s]?|upload[s]?) (say[s]?|mention[s]?|contain[s]?|show[s]?|tell[s]? me)",
            r"what does (my|this|the|that)? ?(document[s]?|file[s]?|note[s]?)? ?say about",
            r"what do (my|this|the|that)? ?(document[s]?|file[s]?|note[s]?)? ?(say|mention)",
            r"(summarize|summary of|give me a summary of) (my|this|the|that)?",
            r"(explain|describe|tell me about) (all |every )?(concepts?|topics?|content|details?|everything|anything)( from| in)?",
            r"from (the )?uploaded",
            r"(can|could) you (please )?tell me",
            r"please (tell|explain|describe|summarize)",
            r"tell me (about|more about|everything about|anything about)?",
            r"i just uploaded (a |the )?document",
            r"find in your memory",
            r"you processed it",
            # Standalone question starters — these are pure filler for embeddings
            r"^what does\b",
            r"^what do\b",
            r"^what is\b",
            r"^what are\b",
        ]

        for pattern in NOISE_PHRASES:
            cleaned = re.sub(pattern, " ", cleaned, flags=re.IGNORECASE)

        # Step 3: Collapse whitespace
        cleaned = re.sub(r"\s+", " ", cleaned).strip()
        cleaned = cleaned.strip(".,?! ")

        # Step 4: If result is empty or only filler, use "summary overview" as anchor
        # This gives the embedding something to work with for broad summarize queries
        _FILLER = {
            "okay", "ok", "please", "can", "you", "could", "would", "just", "now",
            "from", "in", "to", "with", "its", "is", "are", "was", "were", "of", "for",
            "me", "the", "a", "an", "and", "it", "all", "detail", "detailed", "complete",
            "completely", "about", "read", "this", "that", "my", "our",
            # Generic instruction verbs — no semantic signal for retrieval
            "explain", "describe", "summarize", "summary", "everything", "anything",
            "tell", "give", "show", "list", "name", "what", "does", "do", "say",
            "how", "why", "when", "where", "which",
            # Document reference words — useless for embedding search
            "document", "documents", "file", "files", "pdf", "pdfs", "upload", "uploads",
            # Modifier adverbs — no semantic content
            "simply", "more", "briefly", "further", "again", "also",
            # NOTE: Do NOT add domain nouns here (formulas, equations, theorems,
            # concepts, algorithms, JVM, NFA, etc.) — they carry retrieval signal.
        }
        words = cleaned.lower().split()
        meaningful = [w for w in words if w not in _FILLER]

        if not meaningful:
            # Broad summarize query with no topic — use generic anchor
            cleaned = "summary overview key concepts main topics"
            logger.info(f"[ASTRA-CLEAN] No semantic signal — using generic anchor")
        else:
            cleaned = " ".join(meaningful)

        logger.info(f"[ASTRA-CLEAN] '{query[:60]}' → '{cleaned[:60]}'")
        return cleaned

    async def process_file(self, file_path: str) -> Optional[str]:
        """Extract text from various file formats (async-safe)."""
        ext = os.path.splitext(file_path)[1].lower()
        try:
            import asyncio
            if ext == ".pdf":
                return await asyncio.to_thread(self._extract_from_pdf, file_path)
            elif ext == ".docx":
                return await asyncio.to_thread(self._extract_from_docx, file_path)
            elif ext == ".pptx":
                return await asyncio.to_thread(self._extract_from_pptx, file_path)
            elif ext in [".txt", ".md", ".json", ".csv"]:
                return await asyncio.to_thread(self._extract_from_text, file_path)
            elif ext in [".png", ".jpg", ".jpeg"]:
                return await self._extract_from_image(file_path)
            elif ext in [".wav", ".mp3"]:
                return await self._extract_from_audio(file_path)
            else:
                logger.warning(f"Unsupported file format: {ext}")
                return None
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            return None

    def _extract_from_pdf(self, file_path: str) -> str:
        text = ""
        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text

    def _extract_from_docx(self, file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    def _extract_from_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _extract_from_text(self, file_path: str) -> str:
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            if '\x00' in content:
                raise UnicodeError("Null bytes detected, likely UTF-16")
            return content.replace('\ufeff', '').strip()
        except UnicodeError:
            try:
                with open(file_path, "r", encoding="utf-16") as f:
                    content = f.read()
                return content.replace('\x00', '').replace('\ufeff', '').strip()
            except UnicodeError:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read().replace('\x00', '').replace('\ufeff', '').strip()

    async def _extract_from_image(self, file_path: str) -> str:
        """Async image OCR extraction — properly awaits the async vision tool."""
        if not run_vision_ocr:
            return "[Vision Tool Not Loaded]"
        try:
            result = await run_vision_ocr(file_path)
            return result.get("extracted_text", "[No text found in image]")
        except Exception as e:
            logger.error(f"Image OCR failed: {e}")
            return f"[OCR Error: {e}]"

    async def _extract_from_audio(self, file_path: str) -> str:
        """Async audio transcription — properly awaits the async audio tool."""
        if not run_audio_transcription:
            return "[Audio Tool Not Loaded]"
        try:
            result = await run_audio_transcription(file_path)
            return result.get("transcription", "[No transcription found]")
        except Exception as e:
            logger.error(f"Audio transcription failed: {e}")
            return f"[Transcription Error: {e}]"

    async def _get_documents_ordered_by_upload(self, project_id: str) -> List[Dict[str, Any]]:
        """Fetch all documents for a project from SQLite, ordered by upload time."""
        from app.db import engine
        from sqlmodel import Session
        from sqlalchemy import text
        try:
            with Session(engine) as session:
                rows = session.execute(
                    text("SELECT * FROM documents WHERE project_id = :pid ORDER BY uploaded_at ASC"),
                    {"pid": project_id}
                )
                return [dict(r._mapping) for r in rows.fetchall()]
        except Exception as e:
            logger.error(f"Error fetching ordered documents: {e}")
            return []

document_service = DocumentService()
